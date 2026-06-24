#!/usr/bin/env python3
"""
build-lab-slides.py — 2 слайда Лаборатории Модели отрасли БАС
Стиль: Мета-КБ_А2026.pptx — светлый фон, C00000 красный, 4472C4 синий НН, ED7D31 оранжевый Саров
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = sys.argv[1] if len(sys.argv) > 1 else 'data/lab-slides.pptx'

# ── Цвета (из оригинала) ──────────────────────────────────────────────────────
RED    = RGBColor(0xC0, 0x00, 0x00)   # заголовок-метка
DARK   = RGBColor(0x21, 0x21, 0x21)   # основной текст
GREY   = RGBColor(0x70, 0x77, 0x80)   # вторичный текст
GREY2  = RGBColor(0x59, 0x59, 0x59)   # лёгкий серый
BLUE   = RGBColor(0x44, 0x72, 0xC4)   # НН / спецификации
BLUE2  = RGBColor(0x22, 0x4E, 0xA0)   # вторичный синий
ORANGE = RGBColor(0xED, 0x7D, 0x31)   # Саров
ORANGE2= RGBColor(0xE6, 0x5A, 0x00)   # чекпойнт ◆
GREEN  = RGBColor(0x2E, 0x7D, 0x32)   # финальный результат
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LBLUE  = RGBColor(0x1B, 0x39, 0x8A)   # синий для ТЗ/спек
BG     = RGBColor(0xFF, 0xFF, 0xFF)   # белый фон

W = Inches(13.333)
H = Inches(7.5)
PAD = Inches(0.8)

# ── Утилиты ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return slide

def rect(slide, x, y, w, h, fill=None, line_clr=None, line_w=Pt(0.5)):
    s = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_clr:
        s.line.color.rgb = line_clr; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def oval(slide, cx, cy, r, fill):
    s = slide.shapes.add_shape(9, int(cx-r), int(cy-r), int(2*r), int(2*r))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txb(slide, text, x, y, w, h, sz=12, bold=False, clr=DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold; run.font.color.rgb = clr
    return tb

def hline(slide, x, y, w, clr=RGBColor(0xD9,0xD9,0xD9), lw=Pt(0.75)):
    s = slide.shapes.add_connector(1, int(x), int(y), int(x+w), int(y))
    s.line.color.rgb = clr; s.line.width = lw

# ═══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 1 — Лаборатория: Название / Цель / Принципы / Участники
# ═══════════════════════════════════════════════════════════════════════════════
prs = new_prs()
s1  = add_slide(prs)

# Верхняя метка-линия (тонкая красная)
rect(s1, PAD, Inches(0.18), W - PAD*2, Inches(0.04), fill=RED)

# Метка-заголовок (маленький)
txb(s1, 'ЛАБОРАТОРИЯ · АРХИПЕЛАГ 2026', PAD, Inches(0.3), W-PAD*2, Inches(0.25),
    sz=9, bold=False, clr=GREY2)

# Название лаборатории
txb(s1, 'МОДЕЛЬ ОТРАСЛИ БАС', PAD, Inches(0.55), W-PAD*2, Inches(0.5),
    sz=28, bold=True, clr=DARK)

# Подзаголовок
txb(s1, 'Методика исследования отрасли в условиях неполных данных · Исполняемые спецификации как результат',
    PAD, Inches(1.08), W-PAD*2, Inches(0.3),
    sz=12, bold=False, clr=RGBColor(0x45,0x4B,0x54))

# Горизонтальный разделитель
hline(s1, PAD, Inches(1.45), W-PAD*2, clr=RGBColor(0xBD,0xBD,0xBD), lw=Pt(0.75))

# ── Три колонки ─────────────────────────────────────────────────────────────
CX  = [PAD, Inches(5.05), Inches(9.35)]
CW  = Inches(3.9)
TOP = Inches(1.6)
BH  = Inches(5.7)

HEADERS = ['ЦЕЛЬ', 'ПРИНЦИПЫ', 'УЧАСТНИКИ И РОЛИ']
HCOLORS = [RED, RED, RED]
for i, (hdr, hclr) in enumerate(zip(HEADERS, HCOLORS)):
    rect(s1, CX[i], TOP, CW, Inches(0.04), fill=hclr)
    txb(s1, hdr, CX[i], TOP + Inches(0.06), CW, Inches(0.28), sz=11, bold=True, clr=hclr)

# --- ЦЕЛЬ ---
GOALS = [
    ('Методика',     'Воспроизводимый алгоритм исследования отрасли от сбора данных до исполняемых утверждений'),
    ('Параметрическая модель', 'При каком БАС-проникновении поселение становится автономным — доказуемо, не как мнение'),
    ('Исполняемые спецификации', 'Каждый claim запускается как тест: RED / GREEN — единственная защита утверждения'),
    ('Артефакт',     'Отчёт + спека + MCP-инструмент для регулятора, производителя, инвестора'),
]
gy = TOP + Inches(0.42)
for lbl, desc in GOALS:
    rect(s1, CX[0], gy + Inches(0.06), Inches(0.06), Inches(0.22), fill=RED)
    txb(s1, lbl, CX[0]+Inches(0.14), gy, CW-Inches(0.2), Inches(0.26), sz=10, bold=True, clr=DARK)
    txb(s1, desc, CX[0]+Inches(0.14), gy+Inches(0.24), CW-Inches(0.2), Inches(0.52), sz=9, clr=GREY)
    gy += Inches(0.88)

# --- ПРИНЦИПЫ ---
PRINCIPLES = [
    (BLUE, 'Недостаточные данные',
           'Работаем с Росстатом, экспертными оценками, аналогами. Модель итерируется по мере прихода данных — не ждём полной картины.'),
    (ORANGE, 'Факторный анализ',
           'Изолируем один параметр (расстояние, климат, тип сервиса) и смотрим, как он меняет систему. Модель = машина мысленных экспериментов.'),
    (LBLUE, 'Исполняемость',
           'Каждое утверждение — тест. node spec-runner.mjs — единственный судья. Оспорить можно только параметры, не логику.'),
    (GREEN, 'Апофатис честный',
           'Список того, что модель НЕ знает, — часть результата. Это сильнее, чем прятать ограничения.'),
]
py = TOP + Inches(0.42)
for clr, lbl, desc in PRINCIPLES:
    rect(s1, CX[1], py, CW, Inches(0.04), fill=clr)
    txb(s1, lbl, CX[1], py+Inches(0.08), CW, Inches(0.26), sz=10, bold=True, clr=clr)
    txb(s1, desc, CX[1], py+Inches(0.32), CW, Inches(0.58), sz=9, clr=GREY)
    py += Inches(1.0)

# --- УЧАСТНИКИ ---
ROLES = [
    (BLUE,   'ИССЛЕДОВАТЕЛЬ',    'Ставит вопросы, формулирует claims, одобряет граничные условия'),
    (ORANGE, 'ИНЖЕНЕР ДАННЫХ',   'Поставляет данные, калибрует параметры, обновляет data/*.json из реальных источников'),
    (LBLUE,  'АГЕНТ (Мета-КБ)',  'Запускает спеки, анализирует нарушения (RED → разбор), генерирует следующие гипотезы'),
    (GREEN,  'ФАСИЛИТАТОР',      'Ведёт витки: задача → рой → спека → собор → следующий виток'),
]
ry = TOP + Inches(0.42)
for clr, lbl, desc in ROLES:
    rect(s1, CX[2], ry, CW, Inches(0.88), fill=RGBColor(0xF8,0xF8,0xF8), line_clr=clr)
    rect(s1, CX[2], ry, Inches(0.12), Inches(0.88), fill=clr)
    txb(s1, lbl, CX[2]+Inches(0.2), ry+Inches(0.06), CW-Inches(0.28), Inches(0.28),
        sz=9, bold=True, clr=clr)
    txb(s1, desc, CX[2]+Inches(0.2), ry+Inches(0.32), CW-Inches(0.28), Inches(0.52),
        sz=9, clr=GREY)
    ry += Inches(1.05)

# Нижняя подпись
hline(s1, PAD, H-Inches(0.35), W-PAD*2)
txb(s1, 'Лаборатория Модель отрасли БАС · КБ ГаврИИл Код · Архипелаг 2026',
    PAD, H-Inches(0.32), W*0.8, Inches(0.28), sz=8, clr=GREY)

# ═══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 2 — Тайм-план 8 шариков + структура по дням
# ═══════════════════════════════════════════════════════════════════════════════
s2 = add_slide(prs)

# Метка-линия
rect(s2, PAD, Inches(0.18), W-PAD*2, Inches(0.04), fill=RED)
txb(s2, 'ПЛАН РАБОТЫ · СЕМЬ ДНЕЙ + СОЗДАНИЕ ИНСТРУМЕНТОВ', PAD, Inches(0.3), W-PAD*2, Inches(0.28),
    sz=12, bold=True, clr=RED)
txb(s2, 'От задачи до исполняемой спецификации — с образом результата каждого дня',
    PAD, Inches(0.62), W-PAD*2, Inches(0.38), sz=20, bold=True, clr=DARK)

# ── 8 шариков-маркеров (горизонтальная ось) ─────────────────────────────────
BALLS_DATA = [
    ('НН',  BLUE,   '25–27 июл', 'нн'),
    ('Д1',  ORANGE, '3 авг',     'саров'),
    ('Д2',  ORANGE, '4 авг',     'саров'),
    ('Д3',  ORANGE, '5 авг',     'саров'),
    ('Д4',  ORANGE, '6 авг',     'саров'),
    ('Д5',  ORANGE, '7 авг',     'саров'),
    ('Д6',  ORANGE, '8 авг',     'саров'),
    ('ИНС', GREEN,  '9 авг',     'инстр'),
]
NB     = len(BALLS_DATA)
B_Y    = Inches(1.35)
B_R    = Inches(0.26)
B_TOT  = W - PAD*2
B_SPC  = B_TOT / (NB - 1)
B_X0   = PAD

# Горизонтальная линия
hline(s2, B_X0, B_Y, B_TOT, clr=RGBColor(0xBD,0xBD,0xBD), lw=Pt(1.5))

for i, (lbl, clr, day, phase) in enumerate(BALLS_DATA):
    cx = B_X0 + i * B_SPC
    oval(s2, cx, B_Y, B_R, fill=clr)
    txb(s2, lbl, cx - B_R, B_Y - Inches(0.19), B_R*2, Inches(0.38),
        sz=9, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    txb(s2, day, cx - Inches(0.45), B_Y + B_R + Inches(0.04), Inches(0.9), Inches(0.2),
        sz=7, clr=GREY, align=PP_ALIGN.CENTER)

# Подписи фаз
txb(s2, 'НН  25–27 июл', B_X0, B_Y - Inches(0.44), Inches(1.3), Inches(0.22),
    sz=8, bold=True, clr=BLUE)
txb(s2, 'САРОВ  3–8 авг', B_X0 + B_SPC, B_Y - Inches(0.44), Inches(3.8), Inches(0.22),
    sz=8, bold=True, clr=ORANGE)
txb(s2, 'ИНСТРУМЕНТЫ  9 авг', B_X0 + 7*B_SPC - Inches(0.5), B_Y - Inches(0.44), Inches(2.0), Inches(0.22),
    sz=8, bold=True, clr=GREEN)

# ── Таблица Дней (7 строк + строка инструментов) ─────────────────────────────
TBL_TOP = Inches(2.05)
COL_DAY = Inches(0.9)     # ширина колонки "день"
COL_ACT = Inches(4.4)     # "что делает группа"
COL_RES = Inches(6.8)     # "образ результата"
ROW_H   = Inches(0.68)    # высота строки

ROWS = [
    {
        'lbl': 'НН',   'clr': BLUE,
        'action': 'Погрузиться в отрасль: игроки, регуляторика, существующие данные. Сформулировать первый claim.',
        'result': [
            (RED,   '▸ Карта игроков рынка и конкурентов (топ-15)'),
            (RED,   '▸ Реестр регуляторных барьеров (3–5 ключевых)'),
            (LBLUE, '▸ Claim №1 сформулирован: "при dist > 70 км БАС выгоден с первого шага"'),
            (GREEN, '▸ Первый запуск модели: breakeven для 3 поселений-архетипов'),
        ],
        'expanded': True,
    },
    {
        'lbl': 'Д1',   'clr': ORANGE,
        'action': 'Загрузить реальные данные: Росстат, тарифы дронов, ФАП-статистика',
        'result': [(RED, '▸ data/russia-settlements.json обновлён'), (RED, '▸ Базовая линия посчитана для 5 регионов')],
        'expanded': False,
    },
    {
        'lbl': 'Д2',   'clr': ORANGE,
        'action': 'Доказать claim CASCADE@60: структурный скачок при 60% БАС-проникновения',
        'result': [(LBLUE, '▸ bas-industry.spec.mjs :: CASCADE_AT_60 → GREEN (2000 сценариев)'),
                   (RED,   '▸ Контрпримеры проанализированы, граница уточнена')],
        'expanded': False,
    },
    {
        'lbl': 'Д3',   'clr': ORANGE2,
        'action': 'Территориальная модель: где в России внедрять БАС в первую очередь',
        'result': [(ORANGE2,'▸ ◆ Топ-20 приоритетных территорий · checkpoint'),
                   (RED,    '▸ Оптимальное размещение хабов (3 кластера)')],
        'expanded': False,
    },
    {
        'lbl': 'Д4',   'clr': ORANGE,
        'action': 'Системная динамика: 15-летняя симуляция, аттракторы A и B',
        'result': [(LBLUE,'▸ 5 архетипов поселений × 3 сценария регулирования'),
                   (GREEN,'▸ Окно вмешательства: годы 5–8 (максимальный рычаг)')],
        'expanded': False,
    },
    {
        'lbl': 'Д5',   'clr': ORANGE,
        'action': 'Политика: 3 регуляторных сценария, меры поддержки, продуктовая карта дронов',
        'result': [(ORANGE2,'▸ ◆ Разница сценариев: 7 лет · checkpoint'),
                   (RED,    '▸ 3 типа дронов с техническими требованиями')],
        'expanded': False,
    },
    {
        'lbl': 'Д6',   'clr': ORANGE,
        'action': 'Синтез: Модель отрасли БАС — связная картина для регулятора и производителя',
        'result': [(GREEN, '▸ Отраслевой отчёт: утверждения + доказательства + границы'),
                   (GREEN, '▸ Пакет передан Лаб 2 и Закрытой лаб')],
        'expanded': False,
    },
    {
        'lbl': 'ИНС', 'clr': GREEN,
        'action': 'Финализация инструментов: MCP-сервер, исполняемые спеки, документация',
        'result': [(GREEN, '▸ bas-industry MCP — 8 инструментов, GREEN'),
                   (LBLUE,'▸ 4 спеки: CASCADE, BREAKEVEN, TRAP, REGULATORY — сданы')],
        'expanded': False,
    },
]

# Заголовки таблицы
HDR_Y = TBL_TOP - Inches(0.32)
rect(s2, PAD, HDR_Y, COL_DAY, Inches(0.27), fill=RGBColor(0xF2,0xF2,0xF2))
rect(s2, PAD + COL_DAY, HDR_Y, COL_ACT, Inches(0.27), fill=RGBColor(0xF2,0xF2,0xF2))
rect(s2, PAD + COL_DAY + COL_ACT, HDR_Y, COL_RES, Inches(0.27), fill=RGBColor(0xF2,0xF2,0xF2))

for col_txt, col_x, col_w in [
    ('ДЕНЬ', PAD, COL_DAY),
    ('ЧТО ДЕЛАЕТ ЛАБОРАТОРИЯ', PAD+COL_DAY, COL_ACT),
    ('ОБРАЗ РЕЗУЛЬТАТА', PAD+COL_DAY+COL_ACT, COL_RES),
]:
    txb(s2, col_txt, col_x + Inches(0.08), HDR_Y + Inches(0.04), col_w, Inches(0.2),
        sz=11, bold=True, clr=GREY)

# Строки
row_y = TBL_TOP
for row in ROWS:
    clr  = row['clr']
    h    = Inches(1.1) if row['expanded'] else ROW_H

    # Ячейка "день"
    rect(s2, PAD, row_y, COL_DAY, h, fill=clr)
    txb(s2, row['lbl'], PAD, row_y + (h - Inches(0.38))/2, COL_DAY, Inches(0.38),
        sz=16, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)

    # Ячейка "действие"
    rect(s2, PAD+COL_DAY, row_y, COL_ACT, h,
         fill=RGBColor(0xF9,0xF9,0xF9) if not row['expanded'] else RGBColor(0xF0,0xF5,0xFF),
         line_clr=RGBColor(0xE0,0xE0,0xE0))
    txb(s2, row['action'], PAD+COL_DAY+Inches(0.1), row_y+Inches(0.08),
        COL_ACT - Inches(0.2), h - Inches(0.12),
        sz=10, clr=DARK)

    # Ячейка "результат"
    rect(s2, PAD+COL_DAY+COL_ACT, row_y, COL_RES, h,
         line_clr=RGBColor(0xE0,0xE0,0xE0))
    res_x = PAD + COL_DAY + COL_ACT + Inches(0.1)
    res_y_inner = row_y + Inches(0.06)
    for rclr, rtxt in row['result']:
        if res_y_inner + Inches(0.22) > row_y + h - Inches(0.04): break
        txb(s2, rtxt, res_x, res_y_inner, COL_RES - Inches(0.15), Inches(0.24),
            sz=10, bold=True, clr=rclr)
        res_y_inner += Inches(0.26)

    row_y += h + Inches(0.02)

# Нижняя подпись
txb(s2, '◆ контрольная точка — решение принимает человек. Д3: территория, Д5: политика.',
    PAD, row_y + Inches(0.04), W-PAD*2, Inches(0.22), sz=10, bold=True, clr=ORANGE2)

hline(s2, PAD, H-Inches(0.3), W-PAD*2)
txb(s2, 'Лаборатория Модель отрасли БАС · КБ ГаврИИл Код · Архипелаг 2026',
    PAD, H-Inches(0.28), W*0.8, Inches(0.25), sz=8, clr=GREY)

# ── Сохранить ─────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'Saved: {OUT}')
